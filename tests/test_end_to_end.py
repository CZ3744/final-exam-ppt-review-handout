from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from ppt_review_handout.cli_v2 import main


HANDOUT = {
    "chapter_title": "第一章 测试",
    "source_file": "第一章 测试.pptx",
    "review_goals": ["理解测试章节的核心概念。"],
    "knowledge_framework": ["标题", "正文", "表格"],
    "core_points": {"核心概念": ["这是由测试用 handout.json 写入的知识点。"]},
    "terms": {"测试术语": "用于验证渲染链路。"},
    "comparison_tables": [
        {
            "title": "测试对比表",
            "headers": ["项目", "A", "B"],
            "rows": [["特点", "稳定", "可验证"]],
        }
    ],
    "processes": {"测试流程": ["生成 PPTX", "提取 compact.md", "渲染 DOCX"]},
    "exam_points": ["简答：说明测试链路。"],
    "confusing_points": ["区分 slides.json 与 handout.json。"],
    "quick_summary": ["端到端链路应能稳定生成报告和 Word。"],
    "slide_count": 1,
    "image_heavy_slides": [],
}


def make_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "第一章 测试"
    box = slide.shapes.add_textbox(914400, 1371600, 5486400, 914400)
    box.text_frame.text = "这是测试正文。"
    table_shape = slide.shapes.add_table(2, 2, 914400, 2743200, 5486400, 914400)
    table = table_shape.table
    table.cell(0, 0).text = "项目"
    table.cell(0, 1).text = "内容"
    table.cell(1, 0).text = "测试"
    table.cell(1, 1).text = "通过"
    prs.save(path)


def test_extract_then_render_end_to_end(tmp_path: Path):
    ppt_dir = tmp_path / "ppts"
    ppt_dir.mkdir()
    make_pptx(ppt_dir / "第一章 测试.pptx")

    workspace = tmp_path / "workspace"
    assert main(["extract", "--input", str(ppt_dir), "--workspace", str(workspace)]) == 0
    assert (workspace / "extracted" / "第一章 测试.compact.md").exists()
    assert (workspace / "extracted" / "第一章 测试.slides.json").exists()

    analysis = workspace / "analysis"
    analysis.mkdir()
    (analysis / "第一章 测试.handout.json").write_text(json.dumps(HANDOUT, ensure_ascii=False, indent=2), encoding="utf-8")

    output = tmp_path / "outputs"
    assert main(["render", "--analysis", str(analysis), "--output", str(output), "--zip-word", "--layout", "standard"]) == 0
    assert (output / "docx" / "第一章 测试_复习讲义版.docx").exists()
    assert (output / "word_zip" / "review_handouts_docx.zip").exists()
    report = json.loads((output / "report.json").read_text(encoding="utf-8"))
    assert report["docx_count"] == 1
    assert not report["errors"]
