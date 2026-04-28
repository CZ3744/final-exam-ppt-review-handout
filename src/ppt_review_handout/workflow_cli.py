from __future__ import annotations

import argparse
import hashlib
import json
import re
import shutil
import subprocess
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

CN_DIGITS = {"零": 0, "〇": 0, "一": 1, "二": 2, "两": 2, "三": 3, "四": 4, "五": 5, "六": 6, "七": 7, "八": 8, "九": 9}
CN_UNITS = {"十": 10, "百": 100, "千": 1000}
SUPPORTED_SUFFIXES = {".pptx", ".pptm"}
UNSUPPORTED_SUFFIXES = {".ppt"}
LAYOUT_CHOICES = ("review-margin", "standard")
DEFAULT_SECTIONS = {
    "learning_goals": "一、学习/复习目标",
    "outline": "二、内容框架",
    "core_points": "三、核心内容整理",
    "terms": "四、重要概念与术语",
    "comparison_tables": "五、分类、对比与表格",
    "processes": "六、原理、机制与流程",
    "review_focus": "七、重点关注内容",
    "distinctions": "八、易混淆内容辨析",
    "quick_summary": "九、快速总结",
    "image_heavy_slides": "十、需视觉复核的页面",
}
DEFAULT_CONFIG = {
    "language": "zh-CN",
    "output_title_suffix": "",
    "output_filename_suffix": "_handout",
    "zip_filename": "handouts_docx.zip",
    "remove_patterns": [
        "re:^PowerPoint Template$",
        "re:^单击此处编辑母版文本样式$",
        "re:^第二级$",
        "re:^第三级$",
        "re:^第四级$",
        "re:^第五级$",
    ],
    "max_table_rows_in_summary": 8,
    "max_text_items_per_slide": 12,
    "text_clip_limit": 900,
    "image_heavy_threshold": 0.30,
    "docx_font_body": "宋体",
    "docx_font_heading": "黑体",
    "docx_body_size": 10.5,
    "annotation_label": "Notes",
    "report_relative_paths": True,
    "sections": DEFAULT_SECTIONS,
}


def deep_merge(base: dict, patch: dict) -> dict:
    out = dict(base)
    for key, value in (patch or {}).items():
        if isinstance(value, dict) and isinstance(out.get(key), dict):
            out[key] = deep_merge(out[key], value)
        else:
            out[key] = value
    return out


def load_config(path: str | None) -> dict:
    if not path:
        return dict(DEFAULT_CONFIG)
    p = Path(path).expanduser()
    if not p.exists():
        raise FileNotFoundError(f"Config file not found: {p}")
    return deep_merge(DEFAULT_CONFIG, json.loads(p.read_text(encoding="utf-8")))


def chinese_to_int(text: str) -> int:
    text = text.strip()
    if text.isdigit():
        return int(text)
    total = 0
    current = 0
    unit_seen = False
    for ch in text:
        if ch in CN_DIGITS:
            current = CN_DIGITS[ch]
        elif ch in CN_UNITS:
            unit_seen = True
            total += (current or 1) * CN_UNITS[ch]
            current = 0
    total += current
    if not unit_seen and total == 0:
        for ch in text:
            total = total * 10 + CN_DIGITS.get(ch, 0)
    return total


def chapter_index(name: str) -> int:
    match = re.search(r"第\s*([零〇一二两三四五六七八九十百千0-9]+)\s*[章节]", name)
    if match:
        return chinese_to_int(match.group(1))
    match = re.search(r"chapter\s*([0-9]+)", name, re.I)
    if match:
        return int(match.group(1))
    match = re.search(r"(^|[^0-9])([0-9]{1,3})([^0-9]|$)", name)
    return int(match.group(2)) if match else 9999


def discover_pptx(path: Path, recursive: bool = False) -> tuple[list[Path], list[Path]]:
    if not path.exists():
        return [], []
    if path.is_file():
        if path.suffix.lower() in SUPPORTED_SUFFIXES:
            return [path], []
        if path.suffix.lower() in UNSUPPORTED_SUFFIXES:
            return [], [path]
        return [], []
    iterator = path.rglob("*") if recursive else path.iterdir()
    files = [p for p in iterator if p.is_file() and not p.name.startswith("~$")]
    supported = [p for p in files if p.suffix.lower() in SUPPORTED_SUFFIXES]
    unsupported = [p for p in files if p.suffix.lower() in UNSUPPORTED_SUFFIXES]
    return sorted(supported, key=lambda p: (chapter_index(p.name), str(p))), sorted(unsupported, key=lambda p: str(p))


def safe_name(name: str) -> str:
    for ch in '<>:"/\\|?*':
        name = name.replace(ch, "_")
    return " ".join(name.split()).strip() or "untitled"


def stable_id(path: Path) -> str:
    return hashlib.sha1(str(path).encode("utf-8", errors="ignore")).hexdigest()[:8]


def unique_output_stem(path: Path, seen: dict[str, int]) -> str:
    stem = safe_name(path.stem)
    count = seen.get(stem, 0)
    seen[stem] = count + 1
    return stem if count == 0 else f"{stem}_{stable_id(path.resolve())}"


def normalized(text: str) -> str:
    return " ".join(str(text).split()).strip()


def is_noise(text: str, custom: list[str] | None = None) -> bool:
    value = normalized(text)
    if not value:
        return True
    for pattern in custom or []:
        pattern = str(pattern).strip()
        if not pattern:
            continue
        if pattern.startswith("re:") and re.search(pattern[3:], value):
            return True
        if not pattern.startswith("re:") and normalized(pattern) == value:
            return True
    return False


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def extract_shape_text(shape) -> list[str]:
    out: list[str] = []
    if getattr(shape, "has_text_frame", False) and shape.text_frame:
        for para in shape.text_frame.paragraphs:
            run_text = " ".join(run.text for run in para.runs).strip()
            text = run_text or getattr(para, "text", "").strip()
            if text:
                out.append(text)
    return out


def extract_table(shape) -> list[list[str]] | None:
    if not getattr(shape, "has_table", False):
        return None
    return [[normalized(cell.text) for cell in row.cells] for row in shape.table.rows]


def visual_element_count(shape) -> int:
    count = 0
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        count += 1
    if getattr(shape, "has_chart", False):
        count += 1
    if "media" in str(shape.shape_type).lower():
        count += 1
    return count


def extract_notes(slide) -> str:
    try:
        notes = slide.notes_slide.notes_text_frame
        return "\n".join(p.text.strip() for p in notes.paragraphs if p.text.strip())
    except Exception:
        return ""


def detect_role(title: str, texts: list[str], tables: list[list[list[str]]], visuals: int) -> tuple[str, float]:
    joined = " ".join([title] + texts[:6]).lower()
    if tables:
        return "table", 0.70
    if visuals >= 3 and len(texts) <= 5:
        return "figure-heavy", 0.65
    if any(key in joined for key in ["contents", "目录", "agenda"]):
        return "toc", 0.45
    if len(texts) <= 1 and visuals == 0:
        return "transition", 0.40
    return "content", 0.50


def extract_presentation(pptx_path: Path, remove_patterns: list[str] | None = None) -> dict:
    prs = Presentation(str(pptx_path))
    slides = []
    removed = []
    for index, slide in enumerate(prs.slides, start=1):
        texts, tables, visuals = [], [], 0
        for shape in iter_shapes(slide.shapes):
            visuals += visual_element_count(shape)
            table = extract_table(shape)
            if table:
                tables.append(table)
            for text in extract_shape_text(shape):
                if is_noise(text, remove_patterns):
                    removed.append(text)
                else:
                    texts.append(text)
        title = texts[0] if texts and len(texts[0]) <= 80 else f"Slide {index}"
        role, confidence = detect_role(title, texts, tables, visuals)
        slides.append({
            "index": index,
            "title": title,
            "texts": texts[1:] if texts and texts[0] == title else texts,
            "tables": [{"rows": t} for t in tables],
            "visual_element_count": visuals,
            "image_count": visuals,
            "notes": extract_notes(slide),
            "detected_role": role,
            "role_confidence": confidence,
        })
    return {
        "source_file": str(pptx_path),
        "document_title": pptx_path.stem,
        "chapter_title": pptx_path.stem,
        "slide_count": len(slides),
        "slides": slides,
        "removed_boilerplate": sorted(set(removed))[:200],
    }


def clip(text: str, limit: int = 900) -> str:
    text = normalized(text)
    return text if len(text) <= limit else text[: limit - 1] + "…"


def image_heavy_slides(deck: dict) -> list[int]:
    return [s["index"] for s in deck.get("slides", []) if s.get("detected_role") == "figure-heavy"]


def image_heavy_warning(deck: dict, threshold: float) -> str | None:
    total = deck.get("slide_count", 0) or 0
    heavy = image_heavy_slides(deck)
    if total and len(heavy) / total >= threshold:
        return f"Image-heavy deck: {len(heavy)}/{total} slides rely heavily on visuals; review original PPT or run OCR/vision before authoring."
    return None


def deck_to_compact_md(deck: dict, config: dict) -> str:
    rows_limit = max(1, int(config.get("max_table_rows_in_summary", 8)))
    text_limit = max(1, int(config.get("max_text_items_per_slide", 12)))
    char_limit = int(config.get("text_clip_limit", 900))
    lines = [
        f"# {deck.get('document_title') or deck.get('chapter_title')}",
        "",
        f"- Source file: {Path(deck['source_file']).name}",
        f"- Slide count: {deck['slide_count']}",
        "",
        "## Workflow constraints for the calling LLM",
        "",
        "Read this compact file and, when needed, the sibling slides.json. Author the final handout JSON by understanding, merging, and restructuring the source material. Do not mechanically copy slide bullets. Interpret tables semantically. Mark visual-heavy pages for review. Do not invent unsupported conclusions.",
        "",
        "Recommended fields: document_title, source_file, learning_goals, outline, core_points, terms, comparison_tables, processes, review_focus, distinctions, quick_summary, slide_count, image_heavy_slides.",
        "",
        "## Slide summaries",
        "",
    ]
    for slide in deck["slides"]:
        lines.append(f"## Slide {slide['index']}: {slide['title']}")
        lines.append(f"- Detected role: {slide['detected_role']} (confidence {slide.get('role_confidence', 0):.2f})")
        visuals = slide.get("visual_element_count", slide.get("image_count", 0))
        if visuals:
            lines.append(f"- Visual elements: {visuals}; review original slide if the page is mainly visual.")
        if slide.get("texts"):
            lines.append("- Text:")
            for text in slide["texts"][:text_limit]:
                lines.append(f"  - {clip(text, char_limit)}")
            if len(slide["texts"]) > text_limit:
                lines.append(f"  - ... {len(slide['texts']) - text_limit} more items omitted; see slides.json")
        if slide.get("tables"):
            lines.append("- Tables:")
            for i, table in enumerate(slide["tables"], start=1):
                rows = table.get("rows", [])
                lines.append(f"  - Table {i}: {len(rows)} rows")
                for row in rows[:rows_limit]:
                    lines.append("    - " + " | ".join(clip(cell, 80) for cell in row))
                if len(rows) > rows_limit:
                    lines.append(f"    - ... {len(rows) - rows_limit} more rows omitted; see slides.json")
        if slide.get("notes"):
            lines.append("- Notes:")
            lines.append(f"  - {clip(slide['notes'], char_limit)}")
        lines.append("")
    return "\n".join(lines).rstrip() + "\n"


def get_field(handout: dict, primary: str, *aliases, default=None):
    for key in (primary, *aliases):
        if key in handout:
            return handout[key]
    return default


def validate_handout_schema(handout: dict) -> list[str]:
    errors = []
    if "slides" in handout:
        errors.append("This looks like extracted slides.json, not caller-authored handout.json.")
    title = get_field(handout, "document_title", "chapter_title")
    if not isinstance(title, str) or not title.strip():
        errors.append("document_title (or chapter_title) must be a non-empty string.")
    if not isinstance(handout.get("source_file"), str):
        errors.append("source_file must be a string.")
    for names in [("learning_goals", "review_goals"), ("outline", "knowledge_framework"), ("review_focus", "exam_points"), ("distinctions", "confusing_points"), ("quick_summary",), ("image_heavy_slides",)]:
        if not isinstance(get_field(handout, *names, default=[]), list):
            errors.append(f"{names[0]} must be a list.")
    if not isinstance(handout.get("core_points", {}), dict):
        errors.append("core_points must be an object mapping headings to lists.")
    else:
        for key, value in handout.get("core_points", {}).items():
            if not isinstance(value, list):
                errors.append(f"core_points.{key} must be a list.")
    if not isinstance(handout.get("terms", {}), dict):
        errors.append("terms must be an object mapping terms to definitions.")
    if not isinstance(handout.get("comparison_tables", []), list):
        errors.append("comparison_tables must be a list.")
    else:
        for i, table in enumerate(handout.get("comparison_tables", []), start=1):
            if not isinstance(table, dict):
                errors.append(f"comparison_tables[{i}] must be an object.")
                continue
            if not isinstance(table.get("headers", []), list):
                errors.append(f"comparison_tables[{i}].headers must be a list.")
            if not isinstance(table.get("rows", []), list):
                errors.append(f"comparison_tables[{i}].rows must be a list.")
    if not isinstance(handout.get("processes", {}), dict):
        errors.append("processes must be an object mapping process names to step lists.")
    if not isinstance(handout.get("slide_count", 0), int) or handout.get("slide_count", 0) < 0:
        errors.append("slide_count must be a non-negative integer.")
    return errors


def ensure_child(parent, tag: str):
    child = parent.find(qn(tag))
    if child is None:
        child = OxmlElement(tag)
        parent.append(child)
    return child


def set_east_asia_font(element, font_name: str):
    rpr = ensure_child(element, "w:rPr")
    rfonts = ensure_child(rpr, "w:rFonts")
    rfonts.set(qn("w:eastAsia"), font_name)
    rfonts.set(qn("w:hAnsi"), font_name)
    rfonts.set(qn("w:ascii"), font_name)


def set_font(run, font_name="宋体", size=10.5, bold=None):
    run.font.name = font_name
    set_east_asia_font(run._element, font_name)
    run.font.size = Pt(size)
    if bold is not None:
        run.bold = bold


def setup_doc(layout: str, config: dict) -> Document:
    doc = Document()
    sec = doc.sections[0]
    sec.page_width = Cm(21)
    sec.page_height = Cm(29.7)
    margin = 1.5 if layout == "review-margin" else 1.8
    sec.top_margin = sec.bottom_margin = sec.left_margin = sec.right_margin = Cm(margin)
    body = config.get("docx_font_body", "宋体")
    heading = config.get("docx_font_heading", "黑体")
    doc.styles["Normal"].font.size = Pt(float(config.get("docx_body_size", 10.5)))
    set_east_asia_font(doc.styles["Normal"]._element, body)
    for name, size in [("Title", 18), ("Heading 1", 15), ("Heading 2", 13), ("Heading 3", 11)]:
        style = doc.styles[name]
        style.font.size = Pt(size)
        style.font.bold = True
        set_east_asia_font(style._element, heading)
    return doc


def set_cell_border(cell, **edges):
    tc_pr = cell._tc.get_or_add_tcPr()
    borders = tc_pr.find(qn("w:tcBorders"))
    if borders is None:
        borders = OxmlElement("w:tcBorders")
        tc_pr.append(borders)
    for edge, attrs in edges.items():
        tag = "w:" + edge
        element = borders.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            borders.append(element)
        for key, value in attrs.items():
            element.set(qn("w:" + key), str(value))


def prepare_content_target(doc: Document, layout: str, config: dict):
    if layout == "standard":
        return doc
    table = doc.add_table(rows=1, cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    left, right = table.cell(0, 0), table.cell(0, 1)
    left.width = Cm(12.7)
    right.width = Cm(4.2)
    nil = {"val": "nil"}
    for cell in (left, right):
        set_cell_border(cell, top=nil, start=nil, bottom=nil, end=nil)
    set_cell_border(left, end={"val": "single", "sz": "8", "space": "0", "color": "BFBFBF"})
    label = str(config.get("annotation_label", "Notes"))
    if label and right.paragraphs:
        run = right.paragraphs[0].add_run(label)
        set_font(run, config.get("docx_font_body", "宋体"), 8)
    return left


def add_heading(target, text: str, level: int, config: dict):
    if hasattr(target, "add_heading"):
        return target.add_heading(text, level=level)
    p = target.add_paragraph(style=f"Heading {level}")
    r = p.add_run(text)
    set_font(r, config.get("docx_font_heading", "黑体"), 12 if level >= 2 else 14, True)
    return p


def add_bullets(target, items, config: dict, style="List Bullet"):
    for item in items or []:
        if item:
            p = target.add_paragraph(style=style)
            r = p.add_run(str(item))
            set_font(r, config.get("docx_font_body", "宋体"), float(config.get("docx_body_size", 10.5)))


def set_cell_text(cell, text: str, config: dict, bold: bool = False):
    cell.text = ""
    run = cell.paragraphs[0].add_run(str(text))
    set_font(run, config.get("docx_font_heading" if bold else "docx_font_body", "宋体"), float(config.get("docx_body_size", 10.5)), bold)


def add_table(target, title: str, headers: list[str], rows: list[list[str]], layout: str, config: dict):
    if not rows:
        return
    add_heading(target, title or "Table", 3, config)
    width = max(len(headers), max((len(r) for r in rows), default=0), 1)
    if layout == "review-margin" and width >= 4:
        for row in rows:
            p = target.add_paragraph(style="List Bullet")
            r = p.add_run(str(row[0] if row else "Item"))
            set_font(r, config.get("docx_font_heading", "黑体"), 10.5, True)
            for idx, cell in enumerate(row[1:], start=1):
                label = headers[idx] if idx < len(headers) else f"Field {idx+1}"
                p2 = target.add_paragraph()
                p2.paragraph_format.left_indent = Cm(0.6)
                r1 = p2.add_run(f"{label}: ")
                set_font(r1, config.get("docx_font_heading", "黑体"), 10.5, True)
                r2 = p2.add_run(str(cell))
                set_font(r2, config.get("docx_font_body", "宋体"), 10.5)
        return
    table = target.add_table(rows=1, cols=width)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    for i in range(width):
        set_cell_text(table.rows[0].cells[i], headers[i] if i < len(headers) else f"Field {i+1}", config, True)
    for row in rows:
        cells = table.add_row().cells
        for i in range(width):
            set_cell_text(cells[i], row[i] if i < len(row) else "", config)


def handout_to_docx(handout: dict, path: Path, layout: str = "review-margin", config: dict | None = None, title_suffix: str | None = None):
    config = deep_merge(DEFAULT_CONFIG, config or {})
    errors = validate_handout_schema(handout)
    if errors:
        raise ValueError("; ".join(errors))
    doc = setup_doc(layout, config)
    target = prepare_content_target(doc, layout, config)
    title = get_field(handout, "document_title", "chapter_title", default="Untitled")
    suffix = config.get("output_title_suffix", "") if title_suffix is None else title_suffix
    p = target.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(title)
    set_font(r, config.get("docx_font_heading", "黑体"), 18, True)
    if suffix:
        p2 = target.add_paragraph()
        p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r2 = p2.add_run(suffix)
        set_font(r2, config.get("docx_font_heading", "黑体"), 12, True)
    meta = target.add_paragraph()
    set_font(meta.add_run(f"Source: {Path(handout.get('source_file', '')).name}; slides: {handout.get('slide_count', 0)}; layout: {layout}"), config.get("docx_font_body", "宋体"), 9)
    if handout.get("generated_by_fallback"):
        wrn = target.add_paragraph()
        set_font(wrn.add_run("Note: deterministic smoke-test fallback; use only as a rough draft."), config.get("docx_font_heading", "黑体"), 10.5, True)
    sections = config.get("sections", DEFAULT_SECTIONS)
    for key, aliases, style in [
        ("learning_goals", ("review_goals",), "List Bullet"),
        ("outline", ("knowledge_framework",), "List Bullet"),
    ]:
        add_heading(target, sections.get(key, key), 1, config)
        add_bullets(target, get_field(handout, key, *aliases, default=[]), config, style)
    add_heading(target, sections.get("core_points", "Core points"), 1, config)
    for heading, points in (handout.get("core_points") or {}).items():
        add_heading(target, str(heading), 2, config)
        add_bullets(target, points, config)
    add_heading(target, sections.get("terms", "Terms"), 1, config)
    for term, definition in (handout.get("terms") or {}).items():
        p = target.add_paragraph()
        set_font(p.add_run(f"{term}: "), config.get("docx_font_heading", "黑体"), 10.5, True)
        set_font(p.add_run(str(definition)), config.get("docx_font_body", "宋体"), 10.5)
    add_heading(target, sections.get("comparison_tables", "Tables"), 1, config)
    for table in handout.get("comparison_tables") or []:
        add_table(target, table.get("title", "Table"), table.get("headers", []), table.get("rows", []), layout, config)
    add_heading(target, sections.get("processes", "Processes"), 1, config)
    for name, steps in (handout.get("processes") or {}).items():
        add_heading(target, str(name), 2, config)
        add_bullets(target, steps, config, "List Number")
    for key, aliases in [("review_focus", ("exam_points",)), ("distinctions", ("confusing_points",)), ("quick_summary", ())]:
        add_heading(target, sections.get(key, key), 1, config)
        add_bullets(target, get_field(handout, key, *aliases, default=[]), config)
    if handout.get("image_heavy_slides"):
        add_heading(target, sections.get("image_heavy_slides", "Visual review"), 1, config)
        add_bullets(target, [f"Slide {i}: visual-heavy; review original PPT or OCR/vision output." for i in handout["image_heavy_slides"]], config)
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))


def export_pdf(docx_path: Path, pdf_dir: Path) -> tuple[str | None, str | None]:
    exe = shutil.which("libreoffice") or shutil.which("soffice")
    if not exe:
        return None, "PDF export skipped: libreoffice/soffice not found."
    pdf_dir.mkdir(parents=True, exist_ok=True)
    try:
        result = subprocess.run([exe, "--headless", "--convert-to", "pdf", "--outdir", str(pdf_dir), str(docx_path)], check=True, timeout=180, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    except subprocess.TimeoutExpired:
        return None, "PDF export failed: LibreOffice conversion timed out."
    except subprocess.CalledProcessError as exc:
        err = exc.stderr.decode("utf-8", errors="ignore") if exc.stderr else str(exc)
        return None, f"PDF export failed: {err[:500]}"
    pdf = pdf_dir / (docx_path.stem + ".pdf")
    if pdf.exists():
        return str(pdf), None
    out = result.stdout.decode("utf-8", errors="ignore") if result.stdout else ""
    err = result.stderr.decode("utf-8", errors="ignore") if result.stderr else ""
    return None, f"PDF export failed: converted file not found. stdout={out[:200]} stderr={err[:200]}"


def relpath_or_abs(path: str | Path, root: Path, relative: bool = True) -> str:
    p = Path(path)
    if not relative:
        return str(p)
    try:
        return str(p.resolve().relative_to(root.resolve()))
    except Exception:
        return str(p)


def write_report(out_dir: Path, records: list[dict], config: dict | None = None) -> None:
    config = config or DEFAULT_CONFIG
    out_dir.mkdir(parents=True, exist_ok=True)
    relative = bool(config.get("report_relative_paths", True))
    normalized_records = []
    for rec in records:
        item = dict(rec)
        for key in ["source_file", "intermediate_slides", "compact_md", "docx", "pdf"]:
            if item.get(key):
                item[key] = relpath_or_abs(item[key], out_dir, relative)
        normalized_records.append(item)
    summary = {
        "processed_count": len(normalized_records),
        "docx_count": sum(1 for r in normalized_records if r.get("docx")),
        "pdf_count": sum(1 for r in normalized_records if r.get("pdf")),
        "warnings": [w for r in normalized_records for w in r.get("warnings", [])],
        "errors": [e for r in normalized_records for e in r.get("errors", [])],
        "records": normalized_records,
    }
    (out_dir / "report.json").write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    lines = ["# PPT handout workflow report", "", f"- Processed records: {summary['processed_count']}", f"- DOCX files: {summary['docx_count']}", f"- PDF files: {summary['pdf_count']}"]
    if summary["warnings"]:
        lines += ["", "## Warnings"] + [f"- {w}" for w in summary["warnings"]]
    if summary["errors"]:
        lines += ["", "## Errors"] + [f"- {e}" for e in summary["errors"]]
    lines += ["", "## Records"]
    for rec in normalized_records:
        lines.append(f"- {rec.get('document_title') or rec.get('chapter_title')}: {rec.get('source_file')}")
        if rec.get("docx"):
            lines.append(f"  - DOCX: {rec['docx']}")
        if rec.get("pdf"):
            lines.append(f"  - PDF: {rec['pdf']}")
    (out_dir / "report.md").write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")


def extract_cmd(args) -> int:
    try:
        config = load_config(args.config)
    except Exception as exc:
        write_report(Path(args.workspace).expanduser().resolve(), [{"source_file": args.config or "", "document_title": "Config error", "slide_count": 0, "warnings": [], "errors": [str(exc)]}], DEFAULT_CONFIG)
        return 1
    workspace = Path(args.workspace).expanduser().resolve()
    input_path = Path(args.input).expanduser().resolve()
    records = []
    if not input_path.exists():
        records.append({"source_file": str(input_path), "document_title": "Input path not found", "slide_count": 0, "warnings": [], "errors": [f"Input path does not exist: {input_path}"]})
        write_report(workspace, records, config)
        return 1
    files, unsupported = discover_pptx(input_path, recursive=args.recursive)
    for old in unsupported:
        records.append({"source_file": str(old), "document_title": old.stem, "slide_count": 0, "warnings": ["Legacy .ppt is not supported by python-pptx; convert it to .pptx first."], "errors": []})
    if not files:
        if not records:
            records.append({"source_file": str(input_path), "document_title": "No PPTX/PPTM files found", "slide_count": 0, "warnings": [], "errors": ["No supported .pptx/.pptm files found."]})
        write_report(workspace, records, config)
        return 1
    seen = {}
    for ppt in files:
        rec = {"source_file": str(ppt), "document_title": ppt.stem, "slide_count": 0, "warnings": [], "errors": []}
        try:
            deck = extract_presentation(ppt, config.get("remove_patterns"))
            warning = image_heavy_warning(deck, float(config.get("image_heavy_threshold", 0.3)))
            if warning:
                rec["warnings"].append(warning)
            stem = unique_output_stem(ppt, seen)
            out_json = workspace / "extracted" / f"{stem}.slides.json"
            out_md = workspace / "extracted" / f"{stem}.compact.md"
            out_json.parent.mkdir(parents=True, exist_ok=True)
            out_json.write_text(json.dumps(deck, ensure_ascii=False, indent=2), encoding="utf-8")
            out_md.write_text(deck_to_compact_md(deck, config), encoding="utf-8")
            rec.update({"document_title": deck["document_title"], "chapter_title": deck["document_title"], "slide_count": deck["slide_count"], "intermediate_slides": str(out_json), "compact_md": str(out_md)})
        except Exception as exc:
            rec["errors"].append(str(exc))
        records.append(rec)
    write_report(workspace, records, config)
    return 1 if any(r.get("errors") for r in records) else 0


def handout_json_files(analysis: Path, allow_any_json: bool = False) -> list[Path]:
    if analysis.is_file():
        return [analysis] if allow_any_json or analysis.name.endswith(".handout.json") else []
    if not analysis.exists():
        return []
    return sorted(analysis.glob("*.handout.json")) or (sorted(analysis.glob("*.json")) if allow_any_json else [])


def render_cmd(args) -> int:
    try:
        config = load_config(args.config)
    except Exception as exc:
        write_report(Path(args.output).expanduser().resolve(), [{"source_file": args.config or "", "document_title": "Config error", "slide_count": 0, "warnings": [], "errors": [str(exc)]}], DEFAULT_CONFIG)
        return 1
    analysis = Path(args.analysis).expanduser().resolve()
    out = Path(args.output).expanduser().resolve()
    files = handout_json_files(analysis, allow_any_json=args.allow_any_json)
    records = []
    if not files:
        records.append({"source_file": str(analysis), "document_title": "No handout JSON files found", "slide_count": 0, "warnings": [], "errors": ["No *.handout.json files found. Render intentionally refuses raw slides.json files."]})
        write_report(out, records, config)
        return 1
    seen = {}
    for jf in files:
        rec = {"source_file": str(jf), "document_title": jf.stem, "slide_count": 0, "warnings": [], "errors": [], "layout": args.layout}
        try:
            handout = json.loads(jf.read_text(encoding="utf-8"))
            errors = validate_handout_schema(handout)
            if errors:
                raise ValueError("; ".join(errors))
            title = get_field(handout, "document_title", "chapter_title", default=jf.stem)
            base = safe_name(title) + str(config.get("output_filename_suffix", "_handout"))
            count = seen.get(base, 0)
            seen[base] = count + 1
            stem = base if count == 0 else f"{base}_{stable_id(jf.resolve())}"
            docx = out / "docx" / f"{stem}.docx"
            handout_to_docx(handout, docx, layout=args.layout, config=config)
            rec.update({"document_title": title, "chapter_title": title, "slide_count": handout.get("slide_count", 0), "docx": str(docx)})
            if handout.get("generated_by_fallback"):
                rec["warnings"].append("Generated by deterministic smoke-test fallback; semantic LLM authoring is recommended.")
            if args.export_pdf:
                pdf, warning = export_pdf(docx, out / "pdf")
                if pdf:
                    rec["pdf"] = pdf
                if warning:
                    rec["warnings"].append(warning)
        except Exception as exc:
            rec["errors"].append(str(exc))
        records.append(rec)
    if args.zip_word:
        docx_files = [Path(r["docx"]) for r in records if r.get("docx")]
        if docx_files:
            zip_path = out / "word_zip" / str(config.get("zip_filename", "handouts_docx.zip"))
            zip_path.parent.mkdir(parents=True, exist_ok=True)
            with ZipFile(zip_path, "w", ZIP_DEFLATED) as zf:
                for file in docx_files:
                    zf.write(file, arcname=file.name)
    write_report(out, records, config)
    return 1 if any(r.get("errors") for r in records) else 0


def smoke_fallback_handout(deck: dict) -> dict:
    outline, core, terms, tables = [], {}, {}, []
    for slide in deck.get("slides", []):
        title = slide.get("title") or f"Slide {slide.get('index')}"
        if slide.get("detected_role") not in {"title", "transition"} and title not in outline:
            outline.append(title)
        if slide.get("texts"):
            core.setdefault(title, [])
            for text in slide["texts"][:8]:
                core[title].append(text)
                match = re.match(r"^(.{2,24}?)(?:[:：]| is | means | refers to |是指|是|指)(.{8,200})", text, re.I)
                if match and len(terms) < 30:
                    terms[match.group(1).strip()] = match.group(2).strip()
        for i, table in enumerate(slide.get("tables", []), start=1):
            rows = table.get("rows", [])
            if rows:
                tables.append({"title": f"{title} table {i}", "headers": rows[0], "rows": rows[1:]})
    return {
        "document_title": deck.get("document_title") or deck.get("chapter_title"),
        "chapter_title": deck.get("document_title") or deck.get("chapter_title"),
        "source_file": deck.get("source_file", ""),
        "learning_goals": ["Smoke-test draft generated from extracted text; replace with caller-authored semantic goals before final delivery."],
        "review_goals": ["Smoke-test draft generated from extracted text; replace with caller-authored semantic goals before final delivery."],
        "outline": outline[:20],
        "knowledge_framework": outline[:20],
        "core_points": {k: v for k, v in core.items() if v},
        "terms": terms,
        "comparison_tables": tables[:20],
        "processes": {},
        "review_focus": ["Author this section semantically after reading compact.md/slides.json."],
        "exam_points": ["Author this section semantically after reading compact.md/slides.json."],
        "distinctions": [],
        "confusing_points": [],
        "quick_summary": ["This fallback is only for pipeline smoke tests, not a final handout."],
        "slide_count": deck.get("slide_count", 0),
        "image_heavy_slides": image_heavy_slides(deck),
        "generated_by_fallback": True,
    }


def build_cmd(args) -> int:
    if not args.demo_fallback:
        out = Path(args.output).expanduser().resolve()
        write_report(out, [{"source_file": args.input, "document_title": "Build refused", "slide_count": 0, "warnings": [], "errors": ["build requires --demo-fallback. High-quality output must use extract -> caller LLM authors handout.json -> render."]}], DEFAULT_CONFIG)
        return 1
    try:
        config = load_config(args.config)
    except Exception as exc:
        write_report(Path(args.output).expanduser().resolve(), [{"source_file": args.config or "", "document_title": "Config error", "slide_count": 0, "warnings": [], "errors": [str(exc)]}], DEFAULT_CONFIG)
        return 1
    workspace = Path(args.output).expanduser().resolve()
    files, unsupported = discover_pptx(Path(args.input).expanduser().resolve(), recursive=args.recursive)
    records, seen = [], {}
    analysis = workspace / "_fallback_analysis"
    for ppt in files:
        try:
            deck = extract_presentation(ppt, config.get("remove_patterns"))
            stem = unique_output_stem(ppt, seen)
            if args.keep_intermediate:
                inter = workspace / "intermediate"
                inter.mkdir(parents=True, exist_ok=True)
                (inter / f"{stem}.slides.json").write_text(json.dumps(deck, ensure_ascii=False, indent=2), encoding="utf-8")
                (inter / f"{stem}.compact.md").write_text(deck_to_compact_md(deck, config), encoding="utf-8")
            handout = smoke_fallback_handout(deck)
            analysis.mkdir(parents=True, exist_ok=True)
            (analysis / f"{stem}.handout.json").write_text(json.dumps(handout, ensure_ascii=False, indent=2), encoding="utf-8")
        except Exception as exc:
            records.append({"source_file": str(ppt), "document_title": ppt.stem, "slide_count": 0, "warnings": [], "errors": [str(exc)]})
    for old in unsupported:
        records.append({"source_file": str(old), "document_title": old.stem, "slide_count": 0, "warnings": ["Legacy .ppt is not supported by python-pptx; convert it to .pptx first."], "errors": []})
    if records and not files:
        write_report(workspace, records, config)
        return 1 if any(r.get("errors") for r in records) else 0
    render_args = argparse.Namespace(analysis=str(analysis), output=str(workspace), layout=args.layout, export_pdf=args.export_pdf, zip_word=args.zip_word, config=args.config, allow_any_json=False)
    return render_cmd(render_args)


def validate_cmd(args) -> int:
    try:
        load_config(args.config)
    except Exception as exc:
        print(str(exc))
        return 1
    files = handout_json_files(Path(args.analysis).expanduser().resolve(), allow_any_json=args.allow_any_json)
    if not files:
        print("No handout JSON files found.")
        return 1
    ok = True
    for file in files:
        try:
            errors = validate_handout_schema(json.loads(file.read_text(encoding="utf-8")))
        except Exception as exc:
            errors = [str(exc)]
        if errors:
            ok = False
            print(f"[invalid] {file}")
            for err in errors:
                print(f"  - {err}")
        else:
            print(f"[ok] {file}")
    return 0 if ok else 1


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(prog="ppt-review-handout", description="Extract PPTX/PPTM decks, constrain caller-authored handout JSON, and render DOCX/PDF outputs.")
    sub = parser.add_subparsers(dest="command", required=True)
    p = sub.add_parser("extract", help="extract PPTX/PPTM content to slides.json and compact.md")
    p.add_argument("--input", required=True); p.add_argument("--workspace", required=True); p.add_argument("--config"); p.add_argument("--recursive", action="store_true"); p.set_defaults(func=extract_cmd)
    p = sub.add_parser("render", help="render caller-authored *.handout.json files to DOCX/PDF")
    p.add_argument("--analysis", required=True); p.add_argument("--output", required=True); p.add_argument("--config"); p.add_argument("--layout", choices=LAYOUT_CHOICES, default="review-margin"); p.add_argument("--export-pdf", action="store_true"); p.add_argument("--zip-word", action="store_true"); p.add_argument("--allow-any-json", action="store_true"); p.set_defaults(func=render_cmd)
    p = sub.add_parser("validate", help="validate caller-authored handout JSON files")
    p.add_argument("--analysis", required=True); p.add_argument("--config"); p.add_argument("--allow-any-json", action="store_true"); p.set_defaults(func=validate_cmd)
    p = sub.add_parser("build", help="smoke-test fallback pipeline; not a substitute for LLM authoring")
    p.add_argument("--input", required=True); p.add_argument("--output", required=True); p.add_argument("--config"); p.add_argument("--recursive", action="store_true"); p.add_argument("--layout", choices=LAYOUT_CHOICES, default="review-margin"); p.add_argument("--export-pdf", action="store_true"); p.add_argument("--zip-word", action="store_true"); p.add_argument("--keep-intermediate", action="store_true"); p.add_argument("--demo-fallback", action="store_true") ; p.set_defaults(func=build_cmd)
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    return int(args.func(args))


if __name__ == "__main__":
    raise SystemExit(main())
