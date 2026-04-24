from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

CN_DIGITS = {"零": 0, "〇": 0, "一": 1, "二": 2, "两": 2, "三": 3, "四": 4, "五": 5, "六": 6, "七": 7, "八": 8, "九": 9}
CN_UNITS = {"十": 10, "百": 100, "千": 1000}
SUPPORTED_SUFFIXES = {".pptx", ".pptm"}
UNSUPPORTED_SUFFIXES = {".ppt"}
LAYOUT_CHOICES = ("review-margin", "standard")
BOILERPLATE = [
    r"^PowerPoint Template$",
    r"^单击此处编辑母版文本样式$",
    r"^第二级$",
    r"^第三级$",
    r"^第四级$",
    r"^第五级$",
]


def chinese_to_int(text: str) -> int:
    text = text.strip()
    if text.isdigit():
        return int(text)
    total = 0
    current = 0
    unit_seen = False
    for ch in text:
        if ch in CN_DIGITS:
            current = CN_DIGITS[ch]
        elif ch in CN_UNITS:
            unit_seen = True
            unit = CN_UNITS[ch]
            total += (current or 1) * unit
            current = 0
    total += current
    if not unit_seen and total == 0:
        for ch in text:
            total = total * 10 + CN_DIGITS.get(ch, 0)
    return total


def chapter_index(name: str) -> int:
    m = re.search(r"第\s*([零〇一二两三四五六七八九十百千0-9]+)\s*[章节]", name)
    if m:
        return chinese_to_int(m.group(1))
    m = re.search(r"chapter\s*([0-9]+)", name, re.I)
    if m:
        return int(m.group(1))
    m = re.search(r"(^|[^0-9])([0-9]{1,3})([^0-9]|$)", name)
    return int(m.group(2)) if m else 9999


def discover_pptx(path: Path) -> tuple[list[Path], list[Path]]:
    """Return supported PPTX-like files and unsupported legacy PPT files."""
    if path.is_file():
        if path.suffix.lower() in SUPPORTED_SUFFIXES:
            return [path], []
        if path.suffix.lower() in UNSUPPORTED_SUFFIXES:
            return [], [path]
        return [], []
    candidates = [p for p in path.iterdir() if not p.name.startswith("~$")]
    supported = [p for p in candidates if p.suffix.lower() in SUPPORTED_SUFFIXES]
    unsupported = [p for p in candidates if p.suffix.lower() in UNSUPPORTED_SUFFIXES]
    return sorted(supported, key=lambda p: (chapter_index(p.name), p.name)), sorted(unsupported, key=lambda p: p.name)


def safe_name(name: str) -> str:
    for ch in '<>:"/\\|?*':
        name = name.replace(ch, "_")
    return name.strip() or "未命名章节"


def is_noise(text: str, custom: list[str] | None = None) -> bool:
    t = " ".join(str(text).split()).strip()
    if not t:
        return True
    if custom and any(x and x in t for x in custom):
        return True
    return any(re.match(p, t) for p in BOILERPLATE)


def iter_shapes(shapes):
    """Yield shapes recursively, including shapes inside grouped objects."""
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def extract_shape_text(shape) -> list[str]:
    out: list[str] = []
    if getattr(shape, "has_text_frame", False) and shape.text_frame:
        for para in shape.text_frame.paragraphs:
            text = " ".join(run.text for run in para.runs).strip()
            if text:
                out.append(text)
    return out


def extract_table(shape) -> list[list[str]] | None:
    if not getattr(shape, "has_table", False):
        return None
    rows: list[list[str]] = []
    for row in shape.table.rows:
        rows.append([" ".join(cell.text.split()).strip() for cell in row.cells])
    return rows


def extract_notes(slide) -> str:
    try:
        notes = slide.notes_slide.notes_text_frame
        return "\n".join(p.text.strip() for p in notes.paragraphs if p.text.strip())
    except Exception:
        return ""


def detect_role(title: str, texts: list[str], tables: list[list[list[str]]], image_count: int) -> str:
    joined = " ".join([title] + texts[:6])
    if title and len(texts) <= 3 and any(k in joined for k in ["授课人", "课程", "PowerPoint"]):
        return "title"
    if any(k in joined for k in ["主要内容", "目录", "CONTENTS"]):
        return "toc"
    if tables:
        return "table"
    if image_count >= 3 and len(texts) <= 5:
        return "figure-heavy"
    if len(texts) <= 1 and image_count == 0:
        return "transition"
    return "content"


def extract_presentation(pptx_path: Path, remove_patterns: list[str] | None = None) -> dict:
    prs = Presentation(str(pptx_path))
    slides = []
    removed: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        tables: list[list[list[str]]] = []
        image_count = 0
        for shape in iter_shapes(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
            if getattr(shape, "has_table", False):
                table = extract_table(shape)
                if table:
                    tables.append(table)
            for text in extract_shape_text(shape):
                if is_noise(text, remove_patterns):
                    removed.append(text)
                else:
                    texts.append(text)
        title = texts[0] if texts else f"第 {idx} 页"
        if len(title) > 80:
            title = f"第 {idx} 页"
        slides.append(
            {
                "index": idx,
                "title": title,
                "texts": texts[1:] if texts and texts[0] == title else texts,
                "tables": [{"rows": t} for t in tables],
                "image_count": image_count,
                "notes": extract_notes(slide),
                "detected_role": detect_role(title, texts, tables, image_count),
            }
        )
    return {
        "source_file": str(pptx_path),
        "chapter_title": pptx_path.stem,
        "slide_count": len(slides),
        "slides": slides,
        "removed_boilerplate": sorted(set(removed))[:200],
    }


def clip(text: str, limit: int = 900) -> str:
    text = " ".join(str(text).split())
    return text if len(text) <= limit else text[: limit - 1] + "…"


def deck_to_compact_md(deck: dict) -> str:
    lines = [
        f"# {deck['chapter_title']}",
        "",
        f"- 来源文件：{Path(deck['source_file']).name}",
        f"- PPT 页数：{deck['slide_count']}",
        "",
        "## 给调用方 LLM 的任务说明",
        "",
        "请阅读本 compact.md，必要时参考同名 slides.json。你需要自己完成课程内容理解、知识点合并、表格语义化、易考点提炼和复习讲义结构化；不要机械逐页复制 PPT 原文。",
        "",
        "建议输出 handout.json，字段包括：chapter_title, source_file, review_goals, knowledge_framework, core_points, terms, comparison_tables, processes, exam_points, confusing_points, quick_summary, slide_count, image_heavy_slides。",
        "",
        "## 幻灯片摘要",
        "",
    ]
    for slide in deck["slides"]:
        lines.append(f"## 第 {slide['index']} 页：{slide['title']}")
        lines.append(f"- 检测类型：{slide['detected_role']}")
        if slide.get("image_count"):
            lines.append(f"- 图片/图示数量：{slide['image_count']}（如该页主要靠图片表达，请结合原 PPT 复核）")
        if slide.get("texts"):
            lines.append("- 文字内容：")
            for text in slide["texts"][:12]:
                lines.append(f"  - {clip(text)}")
            if len(slide["texts"]) > 12:
                lines.append(f"  - ……其余 {len(slide['texts']) - 12} 条略，详见 slides.json")
        if slide.get("tables"):
            lines.append("- 表格：")
            for i, table in enumerate(slide["tables"], start=1):
                rows = table.get("rows", [])
                lines.append(f"  - 表 {i}：{len(rows)} 行")
                for row in rows[:8]:
                    lines.append("    - " + " | ".join(clip(cell, 80) for cell in row))
                if len(rows) > 8:
                    lines.append(f"    - ……其余 {len(rows) - 8} 行略，详见 slides.json")
        if slide.get("notes"):
            lines.append("- 备注：")
            lines.append(f"  - {clip(slide['notes'])}")
        lines.append("")
    return "\n".join(lines).rstrip() + "\n"


def extract_cmd(args) -> int:
    config = load_config(args.config)
    workspace = Path(args.workspace).expanduser().resolve()
    workspace.mkdir(parents=True, exist_ok=True)
    files, unsupported = discover_pptx(Path(args.input).expanduser().resolve())
    records = []
    for old in unsupported:
        records.append(
            {
                "source_file": str(old),
                "chapter_title": old.stem,
                "slide_count": 0,
                "warnings": ["Legacy .ppt is not supported by python-pptx. Convert it to .pptx first."],
                "errors": [],
            }
        )
    if not files:
        if not records:
            records.append(
                {
                    "source_file": str(args.input),
                    "chapter_title": "No PPTX files found",
                    "slide_count": 0,
                    "warnings": [],
                    "errors": ["No supported .pptx/.pptm files found."],
                }
            )
        write_report(workspace, records)
        return 1
    for ppt in files:
        print(f"[extract] {ppt.name}")
        rec = {"source_file": str(ppt), "chapter_title": ppt.stem, "slide_count": 0, "warnings": [], "errors": []}
        try:
            deck = extract_presentation(ppt, config.get("remove_patterns"))
            stem = safe_name(deck["chapter_title"])
            out_json = workspace / "extracted" / f"{stem}.slides.json"
            out_md = workspace / "extracted" / f"{stem}.compact.md"
            out_json.parent.mkdir(parents=True, exist_ok=True)
            out_json.write_text(json.dumps(deck, ensure_ascii=False, indent=2), encoding="utf-8")
            out_md.write_text(deck_to_compact_md(deck), encoding="utf-8")
            rec.update({"chapter_title": deck["chapter_title"], "slide_count": deck["slide_count"], "intermediate_slides": str(out_json), "compact_md": str(out_md)})
        except Exception as exc:
            rec["errors"].append(str(exc))
        records.append(rec)
    write_report(workspace, records)
    return 1 if any(r["errors"] for r in records) else 0


def load_config(path: str | None) -> dict:
    if not path:
        return {}
    p = Path(path)
    return json.loads(p.read_text(encoding="utf-8")) if p.exists() else {}


def ensure_child(parent, tag: str):
    child = parent.find(qn(tag))
    if child is None:
        child = OxmlElement(tag)
        parent.append(child)
    return child


def ensure_rfonts(element):
    """Ensure `w:rPr/w:rFonts` exists on a run or style XML element."""
    rpr = ensure_child(element, "w:rPr")
    return ensure_child(rpr, "w:rFonts")


def set_east_asia_font(element, font_name: str):
    rfonts = ensure_rfonts(element)
    rfonts.set(qn("w:eastAsia"), font_name)
    rfonts.set(qn("w:hAnsi"), font_name)
    rfonts.set(qn("w:ascii"), font_name)


def set_font(run, east_asia="宋体", size=10.5, bold=None):
    run.font.name = east_asia
    set_east_asia_font(run._element, east_asia)
    run.font.size = Pt(size)
    if bold is not None:
        run.bold = bold


def setup_doc(layout: str = "review-margin") -> Document:
    doc = Document()
    sec = doc.sections[0]
    sec.page_width = Cm(21)
    sec.page_height = Cm(29.7)
    if layout == "review-margin":
        sec.top_margin = Cm(1.5)
        sec.bottom_margin = Cm(1.5)
        sec.left_margin = Cm(1.4)
        sec.right_margin = Cm(1.4)
    else:
        sec.top_margin = sec.bottom_margin = sec.left_margin = sec.right_margin = Cm(1.8)
    styles = doc.styles
    styles["Normal"].font.size = Pt(10.5)
    set_east_asia_font(styles["Normal"]._element, "宋体")
    for name, size in [("Title", 18), ("Heading 1", 15), ("Heading 2", 13), ("Heading 3", 11)]:
        style = styles[name]
        style.font.size = Pt(size)
        style.font.bold = True
        set_east_asia_font(style._element, "黑体")
    return doc


def set_cell_border(cell, **edges):
    """Set cell borders. Example: set_cell_border(cell, end={...})."""
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    tc_borders = tc_pr.find(qn("w:tcBorders"))
    if tc_borders is None:
        tc_borders = OxmlElement("w:tcBorders")
        tc_pr.append(tc_borders)
    for edge in ("top", "start", "bottom", "end", "insideH", "insideV"):
        attrs = edges.get(edge)
        if attrs is None:
            continue
        tag = "w:" + edge
        element = tc_borders.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            tc_borders.append(element)
        for key, value in attrs.items():
            element.set(qn("w:" + key), str(value))


def prepare_content_target(doc: Document, layout: str):
    """Return the object where handout content should be written.

    review-margin uses a two-column invisible table: the left column holds the
    review handout and the right column is intentionally blank for notes. A thin
    vertical border separates the two areas, matching a print-friendly study layout.
    """
    if layout == "standard":
        return doc
    table = doc.add_table(rows=1, cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    try:
        table.autofit = False
    except Exception:
        pass
    left = table.cell(0, 0)
    right = table.cell(0, 1)
    left.width = Cm(12.7)
    right.width = Cm(4.2)
    for row in table.rows:
        row.cells[0].width = Cm(12.7)
        row.cells[1].width = Cm(4.2)
    nil = {"val": "nil"}
    line = {"val": "single", "sz": "8", "space": "0", "color": "BFBFBF"}
    for cell in (left, right):
        set_cell_border(cell, top=nil, start=nil, bottom=nil, end=nil)
    set_cell_border(left, end=line)
    if right.paragraphs:
        p = right.paragraphs[0]
        r = p.add_run("批注区")
        set_font(r, "宋体", 8)
    return left


def add_heading(target, text: str, level: int = 1):
    if hasattr(target, "add_heading"):
        return target.add_heading(text, level=level)
    p = target.add_paragraph(style=f"Heading {level}")
    r = p.add_run(text)
    set_font(r, "黑体", 12 if level >= 2 else 14, True)
    return p


def add_title(target, title: str):
    p = target.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(title)
    set_font(r, "黑体", 18, True)
    p2 = target.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = p2.add_run("考前复习讲义版")
    set_font(r2, "黑体", 12, True)


def add_bullets(target, items, style="List Bullet"):
    for item in items or []:
        if item:
            p = target.add_paragraph(style=style)
            r = p.add_run(str(item))
            set_font(r)


def render_wide_table_as_cards(target, headers: list[str], rows: list[list[str]]):
    for idx, row in enumerate(rows, start=1):
        title = row[0] if row else f"第 {idx} 项"
        p = target.add_paragraph(style="List Bullet")
        r = p.add_run(str(title))
        set_font(r, "黑体", 10.5, True)
        for col_idx, cell in enumerate(row[1:], start=1):
            label = headers[col_idx] if col_idx < len(headers) else f"项目{col_idx + 1}"
            p = target.add_paragraph()
            p.paragraph_format.left_indent = Cm(0.6)
            r1 = p.add_run(f"{label}：")
            set_font(r1, "黑体", 10.5, True)
            r2 = p.add_run(str(cell))
            set_font(r2)


def add_table(target, title: str, headers: list[str], rows: list[list[str]], layout: str = "review-margin"):
    if not rows:
        return
    add_heading(target, title or "对比表", level=3)
    width = max(len(headers), max((len(r) for r in rows), default=0), 1)
    if layout == "review-margin" and width >= 4:
        render_wide_table_as_cards(target, headers, rows)
        return
    table = target.add_table(rows=1, cols=width)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    for i in range(width):
        cell = table.rows[0].cells[i]
        cell.text = headers[i] if i < len(headers) else f"项目{i + 1}"
    for row in rows:
        cells = table.add_row().cells
        for i in range(width):
            cells[i].text = row[i] if i < len(row) else ""


def handout_to_docx(handout: dict, path: Path, layout: str = "review-margin"):
    doc = setup_doc(layout)
    target = prepare_content_target(doc, layout)
    add_title(target, handout.get("chapter_title", "未命名章节"))
    meta = target.add_paragraph()
    r = meta.add_run(f"来源文件：{Path(handout.get('source_file', '')).name}；PPT 页数：{handout.get('slide_count', 0)} 页；版式：{layout}")
    set_font(r, size=9)
    sections = [
        ("一、本章复习目标", handout.get("review_goals", [])),
        ("二、本章知识框架", handout.get("knowledge_framework", [])),
    ]
    for title, items in sections:
        add_heading(target, title, level=1)
        add_bullets(target, items)
    add_heading(target, "三、核心知识点整理", level=1)
    for title, points in (handout.get("core_points") or {}).items():
        add_heading(target, title, level=2)
        add_bullets(target, points)
    add_heading(target, "四、重要概念与名词解释", level=1)
    terms = handout.get("terms") or {}
    if terms:
        for term, definition in terms.items():
            p = target.add_paragraph()
            r1 = p.add_run(f"{term}：")
            set_font(r1, "黑体", 10.5, True)
            r2 = p.add_run(str(definition))
            set_font(r2)
    else:
        add_bullets(target, ["本章未提供单独术语解释，可结合核心知识点复习。"])
    add_heading(target, "五、分类、对比与表格", level=1)
    tables = handout.get("comparison_tables") or []
    if tables:
        for t in tables:
            add_table(target, t.get("title", "对比表"), t.get("headers", []), t.get("rows", []), layout=layout)
    else:
        add_bullets(target, ["本章未提供对比表。"])
    add_heading(target, "六、原理、机制与流程", level=1)
    for name, steps in (handout.get("processes") or {}).items():
        add_heading(target, name, level=2)
        add_bullets(target, steps, style="List Number")
    add_heading(target, "七、易考点归纳", level=1)
    add_bullets(target, handout.get("exam_points", []))
    add_heading(target, "八、易混淆点辨析", level=1)
    add_bullets(target, handout.get("confusing_points", []))
    add_heading(target, "九、本章速记总结", level=1)
    add_bullets(target, handout.get("quick_summary", []))
    if handout.get("image_heavy_slides"):
        add_heading(target, "十、需复核的图片/图示页", level=1)
        add_bullets(target, [f"第 {i} 页图示信息较多，建议结合原 PPT 复核。" for i in handout["image_heavy_slides"]])
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))


def export_pdf(docx_path: Path, pdf_dir: Path) -> tuple[str | None, str | None]:
    exe = shutil.which("libreoffice") or shutil.which("soffice")
    if not exe:
        return None, "PDF export skipped: libreoffice/soffice not found."
    pdf_dir.mkdir(parents=True, exist_ok=True)
    try:
        subprocess.run(
            [exe, "--headless", "--convert-to", "pdf", "--outdir", str(pdf_dir), str(docx_path)],
            check=True,
            timeout=180,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )
    except subprocess.TimeoutExpired:
        return None, "PDF export failed: LibreOffice conversion timed out."
    except subprocess.CalledProcessError as exc:
        err = exc.stderr.decode("utf-8", errors="ignore") if exc.stderr else str(exc)
        return None, f"PDF export failed: {err[:300]}"
    pdf = pdf_dir / (docx_path.stem + ".pdf")
    return (str(pdf), None) if pdf.exists() else (None, "PDF export failed: converted file not found.")


def render_cmd(args) -> int:
    analysis = Path(args.analysis).expanduser().resolve()
    out = Path(args.output).expanduser().resolve()
    files = [analysis] if analysis.is_file() else (sorted(analysis.glob("*.handout.json")) or sorted(analysis.glob("*.json")))
    records = []
    if not files:
        records.append({"source_file": str(analysis), "chapter_title": "No handout JSON files found", "slide_count": 0, "warnings": [], "errors": ["No *.handout.json or *.json files found in analysis directory."]})
        write_report(out, records)
        return 1
    for jf in files:
        rec = {"source_file": str(jf), "chapter_title": jf.stem, "slide_count": 0, "warnings": [], "errors": [], "layout": args.layout}
        try:
            handout = json.loads(jf.read_text(encoding="utf-8"))
            stem = safe_name(handout.get("chapter_title", jf.stem))
            docx = out / "docx" / f"{stem}_复习讲义版.docx"
            handout_to_docx(handout, docx, layout=args.layout)
            rec.update({"chapter_title": handout.get("chapter_title", jf.stem), "slide_count": handout.get("slide_count", 0), "docx": str(docx)})
            if args.export_pdf:
                pdf, warning = export_pdf(docx, out / "pdf")
                if pdf:
                    rec["pdf"] = pdf
                if warning:
                    rec["warnings"].append(warning)
        except Exception as exc:
            rec["errors"].append(str(exc))
        records.append(rec)
    if args.zip_word:
        docx_files = [Path(r["docx"]) for r in records if r.get("docx")]
        if docx_files:
            zp = out / "word_zip" / "review_handouts_docx.zip"
            zp.parent.mkdir(parents=True, exist_ok=True)
            with ZipFile(zp, "w", ZIP_DEFLATED) as zf:
                for f in docx_files:
                    zf.write(f, arcname=f.name)
    write_report(out, records)
    return 1 if any(r["errors"] for r in records) else 0


def fallback_handout(deck: dict) -> dict:
    framework = []
    core: dict[str, list[str]] = {}
    terms: dict[str, str] = {}
    tables = []
    processes: dict[str, list[str]] = {}
    image_heavy = []
    for slide in deck["slides"]:
        title = slide.get("title") or f"第 {slide['index']} 页"
        if slide.get("detected_role") == "figure-heavy":
            image_heavy.append(slide["index"])
        if title not in framework and slide.get("detected_role") not in {"title", "transition", "noise"}:
            framework.append(title)
        if slide.get("texts"):
            core.setdefault(title, [])
            for text in slide["texts"][:8]:
                if len(text) > 8:
                    core[title].append(text)
                m = re.match(r"^(.{2,16}?)(?:是指|是|指)[:：，, ]?(.{8,160})", text)
                if m and len(terms) < 30:
                    terms[m.group(1)] = m.group(2)
        for i, table in enumerate(slide.get("tables", []), start=1):
            rows = table.get("rows", [])
            if rows:
                tables.append({"title": f"{title} 表{i}", "headers": rows[0], "rows": rows[1:]})
        if any(k in title for k in ["流程", "过程", "工艺", "方法", "步骤"]):
            processes[title] = slide.get("texts", [])[:10]
    return {
        "chapter_title": deck["chapter_title"],
        "source_file": deck["source_file"],
        "review_goals": ["掌握本章核心概念、分类、原理、流程与可能考试题型。"],
        "knowledge_framework": framework[:20],
        "core_points": {k: v for k, v in core.items() if v},
        "terms": terms,
        "comparison_tables": tables[:20],
        "processes": {k: v for k, v in processes.items() if v},
        "exam_points": ["围绕本章定义、分类、特点、原理和适用范围准备名词解释与简答题。"],
        "confusing_points": ["对名称相近的材料、工艺或包装技术，应从作用、适用对象和优缺点进行区分。"],
        "quick_summary": ["本章复习时优先掌握框架、概念、分类、对比表和流程。"],
        "slide_count": deck["slide_count"],
        "image_heavy_slides": image_heavy,
    }


def build_cmd(args) -> int:
    workspace = Path(args.output).expanduser().resolve()
    config = load_config(args.config)
    files, unsupported = discover_pptx(Path(args.input).expanduser().resolve())
    records = []
    for old in unsupported:
        records.append({"source_file": str(old), "chapter_title": old.stem, "slide_count": 0, "warnings": ["Legacy .ppt is not supported by python-pptx. Convert it to .pptx first."], "errors": []})
    if not files:
        if not records:
            records.append({"source_file": str(args.input), "chapter_title": "No PPTX files found", "slide_count": 0, "warnings": [], "errors": ["No supported .pptx/.pptm files found."]})
        write_report(workspace, records)
        return 1
    for ppt in files:
        try:
            deck = extract_presentation(ppt, config.get("remove_patterns"))
            if args.keep_intermediate:
                stem = safe_name(deck["chapter_title"])
                inter = workspace / "intermediate"
                inter.mkdir(parents=True, exist_ok=True)
                (inter / f"{stem}.slides.json").write_text(json.dumps(deck, ensure_ascii=False, indent=2), encoding="utf-8")
                (inter / f"{stem}.compact.md").write_text(deck_to_compact_md(deck), encoding="utf-8")
            handout = fallback_handout(deck)
            tmp = workspace / "_fallback_analysis"
            tmp.mkdir(parents=True, exist_ok=True)
            jf = tmp / f"{safe_name(handout['chapter_title'])}.handout.json"
            jf.write_text(json.dumps(handout, ensure_ascii=False, indent=2), encoding="utf-8")
        except Exception as exc:
            records.append({"source_file": str(ppt), "chapter_title": ppt.stem, "slide_count": 0, "warnings": [], "errors": [str(exc)]})
    if records and any(r["errors"] for r in records):
        write_report(workspace, records)
        return 1
    args.analysis = str(workspace / "_fallback_analysis")
    return render_cmd(args)


def write_report(out: Path, records: list[dict]):
    out.mkdir(parents=True, exist_ok=True)
    summary = {
        "input_count": len(records),
        "docx_count": sum(1 for r in records if r.get("docx")),
        "pdf_count": sum(1 for r in records if r.get("pdf")),
        "layouts": sorted(set(r.get("layout", "") for r in records if r.get("layout"))),
        "warnings": [w for r in records for w in r.get("warnings", [])],
        "errors": [e for r in records for e in r.get("errors", [])],
        "records": records,
    }
    (out / "report.json").write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    lines = ["# PPT Review Handout Report", "", f"- Inputs: {summary['input_count']}", f"- DOCX: {summary['docx_count']}", f"- PDF: {summary['pdf_count']}"]
    if summary["layouts"]:
        lines.append(f"- Layout: {', '.join(summary['layouts'])}")
    lines.append("")
    if summary["warnings"]:
        lines += ["## Warnings", ""] + [f"- {w}" for w in summary["warnings"]] + [""]
    if summary["errors"]:
        lines += ["## Errors", ""] + [f"- {e}" for e in summary["errors"]] + [""]
    lines += ["## Records", ""]
    for r in records:
        lines.append(f"- {r.get('chapter_title')}: {r.get('slide_count')} slides")
    (out / "report.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Final-exam PPT review handout skill")
    sub = parser.add_subparsers(dest="command")
    p = sub.add_parser("extract", help="Extract PPTX into slides.json and compact.md")
    p.add_argument("--input", required=True)
    p.add_argument("--workspace", "--output", dest="workspace", required=True)
    p.add_argument("--config")
    p = sub.add_parser("render", help="Render handout.json into Word/PDF")
    p.add_argument("--analysis", required=True)
    p.add_argument("--output", required=True)
    p.add_argument("--export-pdf", action="store_true")
    p.add_argument("--zip-word", action="store_true")
    p.add_argument("--layout", choices=LAYOUT_CHOICES, default="review-margin", help="DOCX layout. review-margin is recommended for study notes and leaves a right-side annotation area.")
    p = sub.add_parser("build", help="One-pass fallback: extract, rule-analyze, render")
    p.add_argument("--input", required=True)
    p.add_argument("--output", required=True)
    p.add_argument("--mode", default="handout")
    p.add_argument("--config")
    p.add_argument("--keep-intermediate", action="store_true")
    p.add_argument("--export-pdf", action="store_true")
    p.add_argument("--zip-word", action="store_true")
    p.add_argument("--layout", choices=LAYOUT_CHOICES, default="review-margin", help="DOCX layout. review-margin is recommended for study notes and leaves a right-side annotation area.")
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    if args.command == "extract":
        return extract_cmd(args)
    if args.command == "render":
        return render_cmd(args)
    if args.command == "build":
        return build_cmd(args)
    parser.print_help()
    return 2


if __name__ == "__main__":
    raise SystemExit(main())
